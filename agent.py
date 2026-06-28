# ────────────────────────────────────────────────────────────────
#  agent.py  –  Multimodal Agentic RAG · Mixture of Experts
#  Ingestion helpers unchanged; ADK MoE layer below.
# ────────────────────────────────────────────────────────────────

import asyncio
import concurrent.futures
import os
import time
import zipfile

import fitz
from docx import Document
from dotenv import load_dotenv
from google import genai
from google.adk.agents import Agent
from google.adk.runners import Runner
from google.adk.sessions import InMemorySessionService
from google.adk.tools.agent_tool import AgentTool
from google.genai import types
from pinecone import Pinecone
from pptx import Presentation

load_dotenv()

# ── API clients ──────────────────────────────────────────────────
client = genai.Client(
    api_key=os.getenv("GOOGLE_API_KEY"),
    http_options=types.HttpOptions(api_version="v1beta"),
)

pc    = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
index = pc.Index("multimodal-rag")

GENERATION_MODEL = "gemini-2.5-flash"
EMBEDDING_MODEL  = "gemini-embedding-001"
EMBEDDING_DIM    = 768


# ════════════════════════════════════════════════════════════════
#  INGESTION HELPERS  (unchanged)
# ════════════════════════════════════════════════════════════════

def chunk_text(text, chunk_size=500, overlap=50):
    words  = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i : i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


def get_embedding(text):
    for attempt in range(3):
        try:
            result = client.models.embed_content(
                model=EMBEDDING_MODEL,
                contents=text,
                config=types.EmbedContentConfig(output_dimensionality=EMBEDDING_DIM),
            )
            return result.embeddings[0].values
        except Exception:
            time.sleep(5 * (attempt + 1))
    return None


def describe_image(image_bytes, ext="png"):
    mime = "image/jpeg" if ext.lower() in ["jpg", "jpeg"] else f"image/{ext.lower()}"
    for attempt in range(3):
        try:
            response = client.models.generate_content(
                model=GENERATION_MODEL,
                contents=[
                    types.Part.from_bytes(data=image_bytes, mime_type=mime),
                    types.Part.from_text(
                        text="Describe this image in detail. If it contains a diagram, "
                             "chart, or flowchart, explain each component and how they connect."
                    ),
                ],
            )
            return response.text
        except Exception as e:
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                return None
            time.sleep(5 * (attempt + 1))
    return None


def extract_text_from_pdf(pdf_path):
    doc, texts = fitz.open(pdf_path), []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            texts.append({"page": page_num + 1, "text": text.strip()})
    return texts


def extract_images_from_pdf(pdf_path):
    doc, images = fitz.open(pdf_path), []
    for page_num in range(len(doc)):
        for img in doc[page_num].get_images():
            base = doc.extract_image(img[0])
            if len(base["image"]) >= 5000:
                images.append({"page": page_num + 1, "data": base["image"], "ext": base["ext"]})
    return images


def extract_from_docx(docx_path):
    doc       = Document(docx_path)
    full_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
    texts     = [{"page": 1, "text": full_text}] if full_text else []
    images    = []
    with zipfile.ZipFile(docx_path) as z:
        for name in z.namelist():
            if name.startswith("word/media/"):
                ext  = name.rsplit(".", 1)[-1].lower()
                data = z.read(name)
                if ext in ["png", "jpg", "jpeg"] and len(data) >= 5000:
                    images.append({"page": 1, "data": data, "ext": ext})
    return texts, images


def extract_from_pptx(pptx_path):
    prs, texts, images = Presentation(pptx_path), [], []
    for slide_num, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"
            try:
                if hasattr(shape, "image"):
                    data = shape.image.blob
                    if len(data) >= 5000:
                        images.append({"page": slide_num + 1, "data": data, "ext": "png"})
            except Exception:
                continue
        if slide_text.strip():
            texts.append({"page": slide_num + 1, "text": slide_text.strip()})
    return texts, images


def ingest_document(file_path):
    ext = file_path.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        texts  = extract_text_from_pdf(file_path)
        images = extract_images_from_pdf(file_path)
    elif ext == "pptx":
        texts, images = extract_from_pptx(file_path)
    elif ext == "docx":
        texts, images = extract_from_docx(file_path)
    else:
        print(f"Unsupported format: {ext}")
        return

    doc_id, batch = os.path.basename(file_path), []

    for item in texts:
        for chunk in chunk_text(item["text"]):
            emb = get_embedding(chunk)
            if emb:
                batch.append({
                    "id":     f"{doc_id}_text_p{item['page']}_{abs(hash(chunk)) % 100000}",
                    "values": emb,
                    "metadata": {
                        "source":  doc_id,
                        "page":    item["page"],
                        "type":    "text",
                        "content": chunk[:1000],
                    },
                })

    for img in images:
        desc = describe_image(img["data"], img["ext"])
        if desc:
            emb = get_embedding(desc)
            if emb:
                batch.append({
                    "id":     f"{doc_id}_img_p{img['page']}_{abs(hash(desc)) % 100000}",
                    "values": emb,
                    "metadata": {
                        "source":  doc_id,
                        "page":    img["page"],
                        "type":    "image",
                        "content": desc[:1000],
                    },
                })

    if batch:
        for i in range(0, len(batch), 100):
            index.upsert(vectors=batch[i : i + 100])
        print(f"✓ {len(batch)} vectors ingested from: {doc_id}")
    else:
        print(f"⚠ No vectors from: {doc_id}")


def batch_ingest(folder_path):
    files = [
        os.path.join(folder_path, f)
        for f in os.listdir(folder_path)
        if f.rsplit(".", 1)[-1].lower() in {"pdf", "pptx", "docx"}
    ]
    print(f"Found {len(files)} document(s)")
    for i, fp in enumerate(files):
        print(f"\n[{i+1}/{len(files)}] {fp}")
        try:
            ingest_document(fp)
        except Exception as e:
            print(f"✗ Failed: {fp} — {e}")
    print("\nDone.")


# ════════════════════════════════════════════════════════════════
#  MoE TOOL FUNCTIONS
# ════════════════════════════════════════════════════════════════

def _pinecone_search(query: str, type_filter: str | None, top_k: int) -> dict:
    """Shared helper — Pinecone query with optional metadata filter."""
    emb = get_embedding(query)
    if not emb:
        return {"error": "Embedding failed — try again shortly."}

    kwargs: dict = dict(vector=emb, top_k=top_k, include_metadata=True)
    if type_filter:
        kwargs["filter"] = {"type": {"$eq": type_filter}}

    results = index.query(**kwargs)
    if not results.matches:
        return {"results": [], "message": "No relevant content found."}

    return {
        "results": [
            {
                "source":  m.metadata.get("source", "unknown"),
                "page":    m.metadata.get("page", "?"),
                "type":    m.metadata.get("type", ""),
                "content": m.metadata.get("content", ""),
                "score":   round(m.score, 3),
            }
            for m in results.matches
        ]
    }


def search_text_chunks(query: str) -> dict:
    """Search only prose/text passages from uploaded documents.

    Use for questions about written content: policies, procedures,
    specifications, project descriptions, and paragraph-form information.

    Args:
        query: A short topic phrase extracted from the request.
               E.g. "Fixxo project name", "PPE policy requirements".

    Returns:
        dict with 'results' list — each has source, page, content, score.
    """
    return _pinecone_search(query, type_filter="text", top_k=5)


def search_image_chunks(query: str) -> dict:
    """Search only image/diagram descriptions from uploaded documents.

    Use for questions about visual elements: charts, flowcharts, figures,
    schematics, diagrams extracted from PDFs, slides, and Word files.

    Args:
        query: Describe the visual element you are looking for.
               E.g. "production flowchart", "safety zone layout diagram".

    Returns:
        dict with 'results' list — each has source, page, content, score.
    """
    return _pinecone_search(query, type_filter="image", top_k=5)


def search_all_chunks(query: str) -> dict:
    """Search ALL content (text + images) with a wider net (top-10).

    Use for cross-document comparisons, broad topic summaries, or questions
    needing both written passages and visual descriptions.

    Args:
        query: Research question or topic. Call 2-3 times with different
               angles for thorough coverage.

    Returns:
        dict with 'results' list (up to 10 items) from text and image chunks.
    """
    return _pinecone_search(query, type_filter=None, top_k=10)


def list_ingested_documents() -> dict:
    """Return all documents currently in the knowledge base.

    Call only when the user explicitly asks what files are available.

    Returns:
        dict with 'documents' (list of filenames) and 'count'.
    """
    try:
        dummy_emb = get_embedding("list all documents")
        if not dummy_emb:
            return {"documents": [], "count": 0}
        results = index.query(vector=dummy_emb, top_k=200, include_metadata=True)
        sources = sorted({
            m.metadata["source"]
            for m in results.matches
            if m.metadata.get("source")
        })
        return {"documents": sources, "count": len(sources)}
    except Exception as e:
        return {"error": str(e), "documents": [], "count": 0}


# ════════════════════════════════════════════════════════════════
#  MIXTURE OF EXPERTS — THREE SPECIALIST SUB-AGENTS
# ════════════════════════════════════════════════════════════════

_CONTEXT_RULE = """
HANDLING CONTEXT FROM PREVIOUS TURNS:
If the request starts with "Context from previous answer:", read that context
carefully before forming your search query — it tells you what the conversation
was about and what entity the user is asking a follow-up question about.
Use it to make your search query specific and self-contained.

Example:
  Request: "Context from previous answer: The previous answer described the
  Fixxo project, an agentic complaint management system. Question: what is
  the project name that you are describing?"
  → Search for: "Fixxo project name" and answer "The project is called Fixxo."

IMPORTANT: NEVER respond with "I am a language model" or "I do not have a
specific project name" without first calling your search tool.
Always search first, then answer based on what the documents say."""


text_expert = Agent(
    name="text_expert",
    model=GENERATION_MODEL,
    description=(
        "Answers questions about written content in uploaded documents — "
        "policies, procedures, project descriptions, specs, and any text-based info."
    ),
    instruction="""You are the Text Expert for Tata Steel SNTI's document knowledge base.

You specialise in text passages: paragraphs, policies, procedures, specifications,
project descriptions, and written explanations from uploaded documents.

RULES (follow strictly every turn):
1. ALWAYS call search_text_chunks before answering. Use a short, precise topic
   query — not the user's full sentence. Re-search with rephrased terms if the
   first attempt returns thin results.
2. Cite every factual claim: [Source: <filename>, Page: <page>]
3. If information is not found, say so. Never fabricate.
4. Use bullet points for multi-item answers; keep prose concise.
5. Do NOT describe images or diagrams — that is the Visual Expert's domain.
""" + _CONTEXT_RULE,
    tools=[search_text_chunks],
)

visual_expert = Agent(
    name="visual_expert",
    model=GENERATION_MODEL,
    description=(
        "Explains diagrams, flowcharts, charts, figures, schematics, and other "
        "visual elements found in uploaded documents."
    ),
    instruction="""You are the Visual Expert for Tata Steel SNTI's document knowledge base.

You specialise in visual content: diagrams, flowcharts, schematics, charts, figures,
and image-based tables. The knowledge base stores Gemini-generated text descriptions
of every image extracted from PDFs, slides, and Word files.

RULES (follow strictly every turn):
1. ALWAYS call search_image_chunks before answering. Describe the visual you are
   looking for (e.g. "blast furnace cross-section"). Try alternative descriptions
   if the first search returns thin results.
2. Walk through each component of a diagram and how they interconnect.
3. Cite every reference: [Source: <filename>, Page/Slide: <page>]
4. If no relevant visual is found, say so. Never guess.
5. Do NOT search written text content — that is the Text Expert's domain.
""" + _CONTEXT_RULE,
    tools=[search_image_chunks],
)

synthesis_expert = Agent(
    name="synthesis_expert",
    model=GENERATION_MODEL,
    description=(
        "Compares, summarises, or analyses topics spanning multiple documents "
        "or requiring both text and visual content."
    ),
    instruction="""You are the Synthesis Expert for Tata Steel SNTI's document knowledge base.

You specialise in cross-document analysis: comparisons across files, broad topic
summaries, trend identification, and questions needing both prose and visuals.

RULES (follow strictly every turn):
1. ALWAYS call search_all_chunks. Run it 2-3 times with different query angles
   to gather comprehensive evidence before composing your answer.
2. Explicitly name the documents and pages your answer draws from.
3. Cite every source: [Source: <filename>, Page: <page>]
4. Use headers to structure multi-part comparisons or analyses.
5. If evidence is insufficient, say what was and wasn't found.
""" + _CONTEXT_RULE,
    tools=[search_all_chunks],
)


# ════════════════════════════════════════════════════════════════
#  ROUTER / ORCHESTRATOR AGENT
# ════════════════════════════════════════════════════════════════

router_agent = Agent(
    name="router",
    model=GENERATION_MODEL,
    description="Mixture-of-Experts orchestrator for Tata Steel SNTI document Q&A.",
    instruction="""You are the Router for an enterprise document Q&A system at Tata Steel SNTI.

Your ONLY job is to classify each query and delegate to the right expert tool.
NEVER answer domain questions yourself.

ROUTING TABLE:
- text_expert       → written content, policies, procedures, project info, specs
- visual_expert     → diagrams, flowcharts, charts, figures, schematics
- synthesis_expert  → cross-document comparison, broad summaries, mixed content
- list_ingested_documents → ONLY when user asks "what files/documents are available?"

══════════════════════════════════════════════════════
FOLLOW-UP QUESTION HANDLING  ← READ THIS CAREFULLY
══════════════════════════════════════════════════════
When the user's question contains words like "this", "that", "it", "the project",
"you mentioned", "above", "its name", "what you described", or any other reference
to the PREVIOUS assistant turn — it is a follow-up question.

For follow-up questions you MUST:
1. Read the most recent assistant message to understand what was being discussed.
2. Build a self-contained, enriched request string:
   "Context from previous answer: [1-2 sentence summary of what was discussed].
    Question: [the user's actual question]"
3. Call the same expert type as before with this enriched request.

Example:
  Previous assistant turn described the Fixxo project.
  User asks: "what is the project name that you are describing"
  → You call text_expert with:
    "Context from previous answer: The previous answer described Fixxo, an
     agentic complaint management system built with NLP and FastAPI.
     Question: What is the name of the project that was described?"

NEVER pass a decontextualized follow-up to an expert — they have no memory
of previous turns and will fail without this context.

══════════════════════════════════════════════════════
OUTPUT FORMAT — EXACTLY THIS, NO EXCEPTIONS:
══════════════════════════════════════════════════════
Return the expert's answer with this format:

<badge>
<blank line>
<expert answer>

Where <badge> is ONE of these lines (choose based on which expert you called):
🔵 Text Expert
🟢 Visual Expert
🟡 Synthesis Expert

The badge MUST be on its own line followed by a blank line.
Do NOT write anything before the badge.
Do NOT add extra commentary of your own after the expert's answer.""",
    tools=[
        AgentTool(agent=text_expert),
        AgentTool(agent=visual_expert),
        AgentTool(agent=synthesis_expert),
        list_ingested_documents,
    ],
)


# ════════════════════════════════════════════════════════════════
#  ADK RUNNER + SESSION SERVICE
# ════════════════════════════════════════════════════════════════

APP_NAME = "enterprise_moe_rag"
USER_ID  = "snti_user"

_session_service = InMemorySessionService()

_runner = Runner(
    app_name=APP_NAME,
    agent=router_agent,
    session_service=_session_service,
)


# ════════════════════════════════════════════════════════════════
#  ASYNC INTERNALS
# ════════════════════════════════════════════════════════════════

async def _create_session_async() -> str:
    session = await _session_service.create_session(
        app_name=APP_NAME,
        user_id=USER_ID,
        state={},
    )
    return session.id


async def _query_agent_async(question: str, session_id: str) -> str:
    content = types.Content(role="user", parts=[types.Part(text=question)])
    async for event in _runner.run_async(
        user_id=USER_ID,
        session_id=session_id,
        new_message=content,
    ):
        if event.is_final_response():
            if event.content and event.content.parts:
                return event.content.parts[0].text
    return "The agent did not produce a response. Please try again."


# ════════════════════════════════════════════════════════════════
#  SYNC WRAPPERS  (Streamlit-compatible)
# ════════════════════════════════════════════════════════════════

_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)


def _run_async(coro):
    """Run a coroutine in a dedicated thread with its own event loop."""
    future = _executor.submit(asyncio.run, coro)
    return future.result()


def create_session() -> str:
    """Create a new ADK session and return its ID."""
    return _run_async(_create_session_async())


def query_agent(question: str, session_id: str) -> str:
    """Route a question through the MoE router and return the final answer.
    Retries up to 3 times with exponential backoff on 503/UNAVAILABLE errors.
    """
    last_error = None
    for attempt in range(3):
        try:
            return _run_async(_query_agent_async(question, session_id))
        except Exception as e:
            last_error = e
            err_str = str(e)
            is_503 = any(x in err_str for x in ["503", "UNAVAILABLE", "high demand"])
            is_429 = any(x in err_str for x in ["429", "RESOURCE_EXHAUSTED", "quota"])
            if is_503 or is_429:
                wait = 8 * (attempt + 1)
                print(f"[503] Model busy, retrying in {wait}s (attempt {attempt+1}/3)")
                time.sleep(wait)
                continue
            raise
    raise RuntimeError(
        "The model is temporarily unavailable after 3 retries. Please try again shortly."
    ) from last_error